from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import pandas as pd
import json
from docx import Document
from pptx import Presentation
import os
import re
import requests
from bs4 import BeautifulSoup
import io
import ollama

def extract_text_from_pdf(file_path):
    pdf = fitz.open(file_path)
    text = []
    for page_num, page in enumerate(pdf):
        page_text = page.get_text()
        if page_text.strip():               # If text found, append it to the list
            text.append(page_text)
        else:                               # If no text found, perform OCR
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            ocr_text = pytesseract.image_to_string(img)
            text.append(ocr_text)
    return "\n".join(text)

def extract_text_from_image(image_path):
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text =[]
    
    for para in doc.paragraphs:
        text.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text.append(cell.text)
    
    rels = doc.part.rels
    for rel in rels:
        rel = rels[rel]
        if "image" in rel.target_ref:
            img_bytes = rel.target_part.blob
            image = Image.open(io.BytesIO(img_bytes))
            ocr_text = pytesseract.image_to_string(image)
            text.append(ocr_text)
    return '\n'.join(text)

def extract_text_from_pptx(file_path):
    ppt = Presentation(file_path)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)
            if shape.shape_type == 13:                      # PICTURE
                image = shape.image
                img_bytes = image.blob
                img = Image.open(io.BytesIO(img_bytes))
                ocr_text = pytesseract.image_to_string(img)
                text.append(ocr_text)
    return '\n'.join(text)

def extract_text_from_xlsx(file_path):
    all_sheets = pd.read_excel(file_path, sheet_name=None)
    text = []
    for sheet_name, df in all_sheets.items():
        text.append(f"Sheet: {sheet_name}")
        text.append(df.to_string(index=False))
    return '\n\n'.join(text)

def extract_text_from_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string(index=False)

def extract_text_from_json(file_path):
    with open(file_path, 'r') as f:
        data = json.load(f)
    return json.dumps(data, indent=4)

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()
    
def extract_text_from_file(file_path):
    end = file_path.split(".")[-1].lower()
    if end == "pdf":
        return extract_text_from_pdf(file_path)
    elif end in ["jpg", "jpeg", "png", "bmp", "tiff"]:
        return extract_text_from_image(file_path)
    elif end == "docx":
        return extract_text_from_docx(file_path)
    elif end == "pptx":
        return extract_text_from_pptx(file_path)
    elif end == "xlsx":
        return extract_text_from_xlsx(file_path)
    elif end == "csv":
        return extract_text_from_csv(file_path)
    elif end == "json":
        return extract_text_from_json(file_path)
    elif end == "txt":
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {end}")

files = ["Mando hackathon PS.pdf"]

def extract_txt_content(files):
    all_text = []
    for file in files:
        text = ""
        text = extract_text_from_file(file)
        all_text.append(text)
    return "\n".join(all_text)

text_content = extract_txt_content(files)

def extract_web_content(text_content):
    urls = re.findall(r'(https?://[^\s]+)', text_content)
    web_content = []
    for url in urls:
        try:
            response = requests.get(url, timeout=5)
            soup = BeautifulSoup(response.content, 'html.parser')
            # Remove scripts/styles
            for script in soup(["script", "style"]):
                script.decompose()
            page_text = soup.get_text(separator='\n')
            page_text = page_text.strip().replace('\n\n', '\n')
            web_content.append(page_text)
        except Exception as e:
            print(f"Failed to fetch {url}: {e}")
    return "\n".join(web_content)

web_content = extract_web_content(text_content)

all_contents = text_content + "\n" + web_content

def chunk_text(text, chunk_size=1000, overlap = 200):
    words = text.split()
    chunks =[]
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(text))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start += chunk_size - overlap
    return chunks

model = SentenceTransformer('all-MiniLM-L6-v2')

chunks = chunk_text(all_contents, chunk_size=100, overlap=20)

embeddings = model.encode(chunks, show_progress_bar=True, convert_to_numpy=True)

index = faiss.IndexFlatL2(embeddings.shape[1])
index.add(embeddings)

def search(query, model, index, chunks, top_k=10):
    query_vec = model.encode([query], convert_to_numpy=True)
    D, I = index.search(query_vec, top_k)
    return [chunks[i] for i in I[0]]

def answer_question(question, contexts, model_name="llama3.2"):
    context = "\n\n".join(contexts)
    prompt = f"""Answer the following question using ONLY the context provided. If the answer isn't present, say "I don't know."

Context:
{context}

Question: {question}

Answer:"""
    response = ollama.generate(model=model_name, prompt=prompt)
    return response['response']

def chat_with_ollama(question, contexts=None, chat_history=None, model_name="llama3.2"):
    """
    Function to chat with Ollama model while maintaining conversation history.
    
    Args:
        question: User's current question
        contexts: Retrieved document chunks relevant to the question
        chat_history: List of previous messages in the conversation
        model_name: Name of the Ollama model to use
        
    Returns:
        response_text: Model's response text
        updated_history: Updated conversation history
    """
    # Initialize chat history if None
    if chat_history is None:
        chat_history = []
    
    # Create system message with context if available
    system_message = "You are a helpful assistant."
    if contexts:
        context_text = "\n\n".join(contexts)
        system_message = f"""You are a helpful assistant. Answer the following question using the context provided. 
        If the answer isn't present in the context, say "I don't know."
        
        Context:
        {context_text}
        """
    
    # If no history exists yet, create a new conversation
    if not chat_history:
        messages = [
            {"role": "system", "content": system_message},
            {"role": "user", "content": question}
        ]
    else:
        # Add the system message with updated context
        messages = [{"role": "system", "content": system_message}]
        # Add the existing conversation history
        messages.extend(chat_history)
        # Add the new user question
        messages.append({"role": "user", "content": question})
    
    # Get response from Ollama
    response = ollama.chat(model=model_name, messages=messages)
    
    # Extract the response content
    response_text = response['message']['content']
    
    # Update chat history with the new user question and model response
    updated_history = chat_history.copy()
    updated_history.append({"role": "user", "content": question})
    updated_history.append({"role": "assistant", "content": response_text})
    
    # Keep only the last 10 messages (5 exchanges) to prevent context window overflow
    if len(updated_history) > 10:
        updated_history = updated_history[-10:]
    
    return response_text, updated_history

def interactive_chat():
    print("Welcome to Document QA Chat! Type 'exit' to end the conversation.")
    print("Type 'new' to start a new conversation or 'upload' to process new files.")
    
    chat_history = []
    current_files = files  # Use your existing files list
    
    while True:
        user_input = input("\nYou: ").strip()
        
        if user_input.lower() == 'exit':
            print("Goodbye!")
            break
            
        elif user_input.lower() == 'new':
            print("Starting a new conversation...")
            chat_history = []
            continue
            
        elif user_input.lower() == 'upload':
            file_path = input("Enter the path to the file you want to upload: ").strip()
            if os.path.exists(file_path):
                current_files.append(file_path)
                # Reprocess all files
                global text_content, web_content, all_contents, chunks, embeddings, index
                text_content = extract_txt_content(current_files)
                web_content = extract_web_content(text_content)
                all_contents = text_content + "\n" + web_content
                chunks = chunk_text(all_contents, chunk_size=100, overlap=20)
                embeddings = model.encode(chunks, show_progress_bar=True, convert_to_numpy=True)
                index = faiss.IndexFlatL2(embeddings.shape[1])
                index.add(embeddings)
                print(f"File {file_path} processed successfully!")
            else:
                print("File not found. Please check the path and try again.")
            continue
        
        # Search for relevant chunks
        relevant_chunks = search(user_input, model, index, chunks, top_k=5)
        
        # Get response and update chat history
        response, chat_history = chat_with_ollama(
            question=user_input,
            contexts=relevant_chunks,
            chat_history=chat_history,
            model_name="llama3.2"
        )
        
        print("\nAssistant:", response)

# Run the interactive chat
if __name__ == "__main__":
    interactive_chat()
