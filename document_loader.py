# document_loader.py

import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import pandas as pd
import json
from docx import Document
from pptx import Presentation
import os
import io

def extract_text_from_pdf(file_path):
    pdf = fitz.open(file_path)
    text = []
    for page in pdf:
        page_text = page.get_text()
        if page_text.strip():
            text.append(page_text)
        else:
            for img in page.get_images(full=True):
                xref = img[0]
                base_image = pdf.extract_image(xref)
                image_bytes = base_image["image"]
                img = Image.open(io.BytesIO(image_bytes))
                ocr_text = pytesseract.image_to_string(img)
                text.append(ocr_text)
    return "\n".join(text)

def extract_text_from_image(image_path):
    image = Image.open(image_path)
    return pytesseract.image_to_string(image)

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = [para.text for para in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text.append(cell.text)
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img = Image.open(io.BytesIO(rel.target_part.blob))
            text.append(pytesseract.image_to_string(img))
    return "\n".join(text)

def extract_text_from_pptx(file_path):
    ppt = Presentation(file_path)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)
            if shape.shape_type == 13:  # Picture
                image = Image.open(io.BytesIO(shape.image.blob))
                text.append(pytesseract.image_to_string(image))
    return "\n".join(text)

def extract_text_from_xlsx(file_path):
    all_sheets = pd.read_excel(file_path, sheet_name=None)
    text = []
    for sheet_name, df in all_sheets.items():
        text.append(f"Sheet: {sheet_name}")
        text.append(df.to_string(index=False))
    return '\n\n'.join(text)

def extract_text_from_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string(index=False)

def extract_text_from_json(file_path):
    with open(file_path, 'r') as f:
        data = json.load(f)
    return json.dumps(data, indent=4)

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_text_from_file(file_path):
    ext = file_path.split('.')[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif ext in ["jpg", "jpeg", "png", "bmp", "tiff"]:
        return extract_text_from_image(file_path)
    elif ext == "docx":
        return extract_text_from_docx(file_path)
    elif ext == "pptx":
        return extract_text_from_pptx(file_path)
    elif ext == "xlsx":
        return extract_text_from_xlsx(file_path)
    elif ext == "csv":
        return extract_text_from_csv(file_path)
    elif ext == "json":
        return extract_text_from_json(file_path)
    elif ext == "txt":
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def extract_txt_content(files):
    all_text = []
    for file in files:
        try:
            text = extract_text_from_file(file)
            all_text.append(text)
        except Exception as e:
            print(f"Error extracting {file}: {e}")
    return "\n".join(all_text)
